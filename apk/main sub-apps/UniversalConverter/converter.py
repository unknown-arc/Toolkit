import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from PIL import Image, ImageTk
import os
import sys
from pathlib import Path
import subprocess
import threading
from datetime import datetime
import webbrowser
import tempfile

# Add FFmpeg to system path dynamically
ffmpeg_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ffmpeg', 'bin')
if os.path.exists(ffmpeg_path):
    os.environ['PATH'] = ffmpeg_path + os.pathsep + os.environ['PATH']

# Check dependencies
def check_dependency(module_name, package_name=None):
    try:
        __import__(module_name)
        return True
    except ImportError:
        print(f"Warning: {module_name} not installed.")
        return False

# Check if FFmpeg is available
def check_ffmpeg():
    try:
        result = subprocess.run(['ffmpeg', '-version'], 
                              capture_output=True, text=True, timeout=5)
        return result.returncode == 0
    except:
        return False

# Import modules if available
FFMPEG_AVAILABLE = check_ffmpeg()
PDF_SUPPORT = check_dependency('PyPDF2') or check_dependency('pypdf2')
DOCX_SUPPORT = check_dependency('docx')
AUDIO_SUPPORT = check_dependency('pydub') and FFMPEG_AVAILABLE
FITZ_SUPPORT = check_dependency('fitz')
REPORTLAB_SUPPORT = check_dependency('reportlab')
PPTX_SUPPORT = check_dependency('pptx')
COMTYPES_SUPPORT = check_dependency('comtypes')

if FITZ_SUPPORT:
    import fitz
if REPORTLAB_SUPPORT:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
if AUDIO_SUPPORT:
    from pydub import AudioSegment
if PPTX_SUPPORT:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

class ModernTooltip:
    def __init__(self, widget, text):
        self.widget = widget
        self.text = text
        self.tip_window = None
        self.widget.bind("<Enter>", self.show_tip)
        self.widget.bind("<Leave>", self.hide_tip)

    def show_tip(self, event=None):
        if self.tip_window:
            return
        x, y, _, _ = self.widget.bbox("insert")
        x += self.widget.winfo_rootx() + 25
        y += self.widget.winfo_rooty() + 25
        
        self.tip_window = tw = tk.Toplevel(self.widget)
        tw.wm_overrideredirect(True)
        tw.wm_geometry(f"+{x}+{y}")
        
        label = tk.Label(tw, text=self.text, justify=tk.LEFT,
                        background="#ffffe0", relief=tk.SOLID, borderwidth=1,
                        font=("Arial", 10))
        label.pack()

    def hide_tip(self, event=None):
        if self.tip_window:
            self.tip_window.destroy()
            self.tip_window = None

class UniversalConverter:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title("Universal File Converter Pro v2.3")
        self.root.geometry("800x700")
        self.root.minsize(750, 650)
        
        # Initialize conversion statistics FIRST
        self.conversion_stats = {
            'total_files': 0,
            'successful': 0,
            'failed': 0,
            'start_time': None
        }
        
        # Configure style
        self.setup_styles()
        self.setup_file_types()
        self.setup_ui()
        
        # Show FFmpeg status
        if not FFMPEG_AVAILABLE:
            self.log_message("⚠️ FFmpeg not found - audio/video conversion disabled")
            self.log_message("   Download FFmpeg and place in ffmpeg/bin/ folder")
        
    def setup_styles(self):
        self.style = ttk.Style()
        self.style.configure("TButton", padding=6, relief="flat", background="#4CAF50")
        self.style.configure("Title.TLabel", font=("Arial", 16, "bold"), foreground="#2C3E50")
        self.style.configure("Subtitle.TLabel", font=("Arial", 12), foreground="#34495E")
        
    def setup_file_types(self):
        self.supported_formats = {
            'Images': {
                'extensions': ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp', '.tiff', '.ico'],
                'convert_to': ['PNG', 'JPEG', 'BMP', 'GIF', 'WEBP', 'PDF', 'PPTX', 'TIFF'],
                'category': 'image',
                'icon': '🖼️'
            },
            'PDF': {
                'extensions': ['.pdf'],
                'convert_to': ['PNG', 'JPG', 'TXT', 'PPTX'] + (['PDF/A'] if FITZ_SUPPORT else []),
                'category': 'document',
                'icon': '📄'
            },
            'Documents': {
                'extensions': ['.docx', '.doc', '.txt', '.rtf'],
                'convert_to': ['PDF', 'TXT'] + (['DOCX'] if DOCX_SUPPORT else []),
                'category': 'document',
                'icon': '📝'
            },
            'Presentations': {
                'extensions': ['.pptx', '.ppt'],
                'convert_to': ['PDF', 'PNG', 'JPG', 'TXT'] + (['PPTX'] if PPTX_SUPPORT else []),
                'category': 'presentation',
                'icon': '📊'
            },
            'Audio': {
                'extensions': ['.mp3', '.wav', '.ogg', '.flac', '.m4a', '.aac'],
                'convert_to': ['MP3', 'WAV', 'OGG', 'FLAC'] if AUDIO_SUPPORT else [],
                'category': 'audio',
                'icon': '🎵'
            }
        }
        
        # Remove empty categories
        self.supported_formats = {k: v for k, v in self.supported_formats.items() if v['convert_to']}
        
    def setup_ui(self):
        # Main container with modern look
        main_container = tk.Frame(self.root, bg='#f8f9fa')
        main_container.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # Header with gradient effect
        header_frame = tk.Frame(main_container, bg='#2C3E50', height=80)
        header_frame.pack(fill=tk.X, pady=(0, 10))
        header_frame.pack_propagate(False)
        
        title_label = tk.Label(header_frame, text="Universal File Converter Pro v2.3", 
                              font=("Arial", 20, "bold"), fg="white", bg='#2C3E50')
        title_label.pack(pady=20)
        
        subtitle_label = tk.Label(header_frame, text="Convert Any File Format with Ease", 
                                font=("Arial", 10), fg="#BDC3C7", bg='#2C3E50')
        subtitle_label.pack(pady=(0, 10))
        
        # Quick actions frame
        quick_actions = tk.Frame(main_container, bg='#f8f9fa')
        quick_actions.pack(fill=tk.X, pady=5)
        
        tk.Button(quick_actions, text="📁 Quick Image to PDF", command=self.quick_image_to_pdf,
                 bg="#E74C3C", fg="white", font=("Arial", 9), width=18).pack(side=tk.LEFT, padx=2)
        tk.Button(quick_actions, text="🎵 Audio Batch Convert", command=self.quick_audio_batch,
                 bg="#9B59B6", fg="white", font=("Arial", 9), width=18).pack(side=tk.LEFT, padx=2)
        tk.Button(quick_actions, text="📊 Create Slideshow", command=self.quick_slideshow,
                 bg="#3498DB", fg="white", font=("Arial", 9), width=18).pack(side=tk.LEFT, padx=2)
        tk.Button(quick_actions, text="🔄 PDF to PPT", command=self.quick_pdf_to_ppt,
                 bg="#27AE60", fg="white", font=("Arial", 9), width=15).pack(side=tk.LEFT, padx=2)
        
        # Dependency status
        deps_frame = tk.Frame(main_container, relief=tk.GROOVE, bd=1, bg='white')
        deps_frame.pack(fill=tk.X, pady=5)
        
        deps_text = "Available features: "
        deps_list = []
        if PDF_SUPPORT or FITZ_SUPPORT:
            deps_list.append("PDF")
        if DOCX_SUPPORT:
            deps_list.append("DOCX")
        if PPTX_SUPPORT:
            deps_list.append("PPT")
        if AUDIO_SUPPORT:
            deps_list.append("Audio")
        if REPORTLAB_SUPPORT:
            deps_list.append("Reports")
            
        deps_text += ", ".join(deps_list) if deps_list else "Basic image conversion only"
        
        deps_label = tk.Label(deps_frame, text=deps_text, font=("Arial", 9), fg="green", bg='white')
        deps_label.pack(pady=3)
        
        # FFmpeg status
        if not FFMPEG_AVAILABLE:
            ffmpeg_frame = tk.Frame(main_container, relief=tk.GROOVE, bd=1, bg="#FFF3CD")
            ffmpeg_frame.pack(fill=tk.X, pady=2)
            ffmpeg_label = tk.Label(ffmpeg_frame, 
                                   text="⚠️ FFmpeg not found - audio conversion disabled", 
                                   font=("Arial", 8), fg="#856404", bg="#FFF3CD")
            ffmpeg_label.pack(pady=2)
        
        # File selection
        file_frame = tk.LabelFrame(main_container, text=" File Selection ", font=("Arial", 11, "bold"),
                                  bg='#f8f9fa', fg="#2C3E50", padx=10, pady=10)
        file_frame.pack(fill=tk.X, pady=10)
        
        file_input_frame = tk.Frame(file_frame, bg='#f8f9fa')
        file_input_frame.pack(fill=tk.X)
        
        tk.Label(file_input_frame, text="Select File/Folder:", font=("Arial", 10), 
                bg='#f8f9fa').pack(side=tk.LEFT)
        self.file_path = tk.StringVar()
        file_entry = tk.Entry(file_input_frame, textvariable=self.file_path, width=50, 
                             font=("Arial", 10), relief=tk.SOLID, bd=1)
        file_entry.pack(side=tk.LEFT, padx=5, fill=tk.X, expand=True)
        
        button_frame = tk.Frame(file_input_frame, bg='#f8f9fa')
        button_frame.pack(side=tk.LEFT)
        
        tk.Button(button_frame, text="📁 Browse File", command=self.browse_file, 
                 bg="#4CAF50", fg="white", font=("Arial", 9), width=12).pack(side=tk.LEFT, padx=2)
        tk.Button(button_frame, text="📂 Browse Folder", command=self.browse_folder, 
                 bg="#2196F3", fg="white", font=("Arial", 9), width=12).pack(side=tk.LEFT, padx=2)
        tk.Button(button_frame, text="🗑️ Clear", command=self.clear_selection, 
                 bg="#95a5a6", fg="white", font=("Arial", 9), width=8).pack(side=tk.LEFT, padx=2)
        
        # Conversion settings frame
        settings_frame = tk.LabelFrame(main_container, text=" Conversion Settings ", 
                                      font=("Arial", 11, "bold"), bg='#f8f9fa', fg="#2C3E50", 
                                      padx=10, pady=10)
        settings_frame.pack(fill=tk.X, pady=10)
        
        # Format and quality in one row
        format_quality_frame = tk.Frame(settings_frame, bg='#f8f9fa')
        format_quality_frame.pack(fill=tk.X, pady=5)
        
        tk.Label(format_quality_frame, text="Convert to:", font=("Arial", 10), 
                bg='#f8f9fa').pack(side=tk.LEFT)
        self.target_format = tk.StringVar(value="PNG")
        self.format_combo = ttk.Combobox(format_quality_frame, textvariable=self.target_format, 
                                       state="readonly", width=15)
        self.format_combo.pack(side=tk.LEFT, padx=5)
        
        # Quality only for image formats
        self.quality_label = tk.Label(format_quality_frame, text="Quality:", font=("Arial", 10), 
                                     bg='#f8f9fa')
        self.quality_label.pack(side=tk.LEFT, padx=(20, 0))
        self.quality = tk.IntVar(value=85)
        self.quality_scale = tk.Scale(format_quality_frame, from_=1, to=100, variable=self.quality, 
                                     orient=tk.HORIZONTAL, length=120, showvalue=True, 
                                     bg='#f8f9fa', highlightbackground='#f8f9fa')
        self.quality_scale.pack(side=tk.LEFT, padx=5)
        
        # Advanced options
        advanced_frame = tk.Frame(settings_frame, bg='#f8f9fa')
        advanced_frame.pack(fill=tk.X, pady=5)
        
        # PPT Layout (only show when converting to PPTX)
        self.ppt_layout_label = tk.Label(advanced_frame, text="PPT Layout:", font=("Arial", 10), 
                                        bg='#f8f9fa')
        self.ppt_layout_label.pack(side=tk.LEFT)
        self.ppt_layout = tk.StringVar(value="Single Image")
        self.ppt_layout_combo = ttk.Combobox(advanced_frame, textvariable=self.ppt_layout, 
                                           values=["Single Image", "Multiple Images - One per Slide", 
                                                  "Grid Layout (2x2)", "Grid Layout (3x3)"],
                                           state="readonly", width=25)
        self.ppt_layout_combo.pack(side=tk.LEFT, padx=5)
        
        # PDF to PPT options
        self.pdf_ppt_options_label = tk.Label(advanced_frame, text="PDF to PPT:", font=("Arial", 10), 
                                            bg='#f8f9fa')
        self.pdf_ppt_options = tk.StringVar(value="One slide per page")
        self.pdf_ppt_options_combo = ttk.Combobox(advanced_frame, textvariable=self.pdf_ppt_options,
                                                values=["One slide per page", "Extract text content", 
                                                       "High quality images", "Fast conversion"],
                                                state="readonly", width=25)
        
        # Output options
        tk.Label(advanced_frame, text="Output:", font=("Arial", 10), 
                bg='#f8f9fa').pack(side=tk.LEFT, padx=(20, 0))
        self.output_option = tk.StringVar(value="Same folder")
        output_combo = ttk.Combobox(advanced_frame, textvariable=self.output_option,
                                  values=["Same folder", "Custom folder", "Desktop"],
                                  state="readonly", width=12)
        output_combo.pack(side=tk.LEFT, padx=5)
        
        # Initially hide options
        self.toggle_ppt_options(False)
        self.toggle_pdf_ppt_options(False)
        self.toggle_quality_option(True)
        
        # Buttons frame
        button_frame = tk.Frame(main_container, bg='#f8f9fa')
        button_frame.pack(pady=15)
        
        convert_btn = tk.Button(button_frame, text="🔄 Convert Single File", 
                 command=self.convert_file, 
                 bg="#2196F3", fg="white", font=("Arial", 12, "bold"),
                 width=18, height=2, cursor="hand2")
        convert_btn.pack(side=tk.LEFT, padx=5)
        ModernTooltip(convert_btn, "Convert the selected file to the target format")
        
        batch_btn = tk.Button(button_frame, text="📁 Batch Convert", 
                 command=self.batch_convert, 
                 bg="#FF9800", fg="white", font=("Arial", 12, "bold"),
                 width=18, height=2, cursor="hand2")
        batch_btn.pack(side=tk.LEFT, padx=5)
        ModernTooltip(batch_btn, "Convert all supported files in the selected folder")
        
        settings_btn = tk.Button(button_frame, text="⚙️ Settings", 
                 command=self.show_settings,
                 bg="#607D8B", fg="white", font=("Arial", 12, "bold"),
                 width=12, height=2, cursor="hand2")
        settings_btn.pack(side=tk.LEFT, padx=5)
        ModernTooltip(settings_btn, "Configure converter settings")
        
        # Progress bar with percentage
        progress_frame = tk.Frame(main_container, bg='#f8f9fa')
        progress_frame.pack(fill=tk.X, pady=10)
        
        self.progress = ttk.Progressbar(progress_frame, mode='determinate')
        self.progress.pack(fill=tk.X, side=tk.LEFT, expand=True)
        
        self.progress_label = tk.Label(progress_frame, text="0%", font=("Arial", 9), 
                                      bg='#f8f9fa', fg="#2C3E50")
        self.progress_label.pack(side=tk.RIGHT, padx=5)
        
        # Status with time
        status_frame = tk.Frame(main_container, bg='#f8f9fa')
        status_frame.pack(fill=tk.X, pady=5)
        
        self.status = tk.Label(status_frame, text="Ready to convert files", 
                              font=("Arial", 10), fg="gray", bg='#f8f9fa')
        self.status.pack(side=tk.LEFT)
        
        self.time_label = tk.Label(status_frame, text="", font=("Arial", 9), 
                                  fg="#7F8C8D", bg='#f8f9fa')
        self.time_label.pack(side=tk.RIGHT)
        
        # Notebook for info and logs
        notebook = ttk.Notebook(main_container)
        notebook.pack(fill=tk.BOTH, expand=True, pady=10)
        
        # File info tab
        info_frame = tk.Frame(notebook, bg='white')
        notebook.add(info_frame, text="File Information")
        
        self.info_text = tk.Text(info_frame, height=10, width=70, font=("Arial", 9),
                                wrap=tk.WORD, relief=tk.FLAT, bg='#f8f9fa')
        scrollbar_info = tk.Scrollbar(info_frame, command=self.info_text.yview)
        self.info_text.config(yscrollcommand=scrollbar_info.set)
        self.info_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=5, pady=5)
        scrollbar_info.pack(side=tk.RIGHT, fill=tk.Y, pady=5)
        self.info_text.insert(tk.END, "No file selected")
        self.info_text.config(state=tk.DISABLED)
        
        # Conversion log tab
        log_frame = tk.Frame(notebook, bg='white')
        notebook.add(log_frame, text="Conversion Log")
        
        self.log_text = tk.Text(log_frame, height=10, width=70, font=("Arial", 8),
                               wrap=tk.WORD, relief=tk.FLAT, bg='#f8f9fa')
        scrollbar_log = tk.Scrollbar(log_frame, command=self.log_text.yview)
        self.log_text.config(yscrollcommand=scrollbar_log.set)
        self.log_text.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=5, pady=5)
        scrollbar_log.pack(side=tk.RIGHT, fill=tk.Y, pady=5)
        self.log_text.config(state=tk.DISABLED)
        
        # Statistics tab
        stats_frame = tk.Frame(notebook, bg='white')
        notebook.add(stats_frame, text="Statistics")
        
        self.stats_text = tk.Text(stats_frame, height=10, width=70, font=("Arial", 9),
                                 wrap=tk.WORD, relief=tk.FLAT, bg='#f8f9fa')
        self.stats_text.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        self.update_stats()
        self.stats_text.config(state=tk.DISABLED)
        
        # Update time periodically
        self.update_time()
    
    def toggle_ppt_options(self, show=True):
        if show:
            self.ppt_layout_label.pack(side=tk.LEFT)
            self.ppt_layout_combo.pack(side=tk.LEFT, padx=5)
        else:
            self.ppt_layout_label.pack_forget()
            self.ppt_layout_combo.pack_forget()
    
    def toggle_pdf_ppt_options(self, show=True):
        if show:
            self.pdf_ppt_options_label.pack(side=tk.LEFT)
            self.pdf_ppt_options_combo.pack(side=tk.LEFT, padx=5)
        else:
            self.pdf_ppt_options_label.pack_forget()
            self.pdf_ppt_options_combo.pack_forget()
    
    def toggle_quality_option(self, show=True):
        if show:
            self.quality_label.pack(side=tk.LEFT, padx=(20, 0))
            self.quality_scale.pack(side=tk.LEFT, padx=5)
        else:
            self.quality_label.pack_forget()
            self.quality_scale.pack_forget()
    
    def update_time(self):
        current_time = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        self.time_label.config(text=current_time)
        self.root.after(1000, self.update_time)
    
    def clear_selection(self):
        self.file_path.set("")
        self.info_text.config(state=tk.NORMAL)
        self.info_text.delete(1.0, tk.END)
        self.info_text.insert(tk.END, "No file selected")
        self.info_text.config(state=tk.DISABLED)
        self.status.config(text="Selection cleared")
    
    def show_settings(self):
        settings_window = tk.Toplevel(self.root)
        settings_window.title("Converter Settings")
        settings_window.geometry("400x400")
        settings_window.transient(self.root)
        settings_window.grab_set()
        
        tk.Label(settings_window, text="Converter Settings", font=("Arial", 16, "bold")).pack(pady=10)
        
        # Settings frame
        settings_content = tk.Frame(settings_window)
        settings_content.pack(fill=tk.BOTH, expand=True, padx=20, pady=10)
        
        # Default output format
        tk.Label(settings_content, text="Default Output Format:", font=("Arial", 10)).pack(anchor=tk.W)
        default_format = ttk.Combobox(settings_content, values=["PNG", "JPEG", "PDF", "MP3", "PPTX"], state="readonly")
        default_format.set("PNG")
        default_format.pack(fill=tk.X, pady=5)
        
        # Default quality
        tk.Label(settings_content, text="Default Image Quality:", font=("Arial", 10)).pack(anchor=tk.W)
        quality_frame = tk.Frame(settings_content)
        quality_frame.pack(fill=tk.X, pady=5)
        quality_var = tk.IntVar(value=85)
        quality_scale = tk.Scale(quality_frame, from_=1, to=100, variable=quality_var, 
                                orient=tk.HORIZONTAL, showvalue=True)
        quality_scale.pack(fill=tk.X)
        
        # Overwrite existing files
        overwrite_var = tk.BooleanVar(value=False)
        tk.Checkbutton(settings_content, text="Overwrite existing files without warning", 
                      variable=overwrite_var, font=("Arial", 10)).pack(anchor=tk.W, pady=5)
        
        # Create subfolders for batch operations
        subfolder_var = tk.BooleanVar(value=True)
        tk.Checkbutton(settings_content, text="Create subfolders for batch conversions", 
                      variable=subfolder_var, font=("Arial", 10)).pack(anchor=tk.W, pady=5)
        
        # Button frame
        button_frame = tk.Frame(settings_content)
        button_frame.pack(fill=tk.X, pady=20)
        
        tk.Button(button_frame, text="Save Settings", command=settings_window.destroy,
                 bg="#4CAF50", fg="white", width=15).pack(side=tk.LEFT, padx=5)
        tk.Button(button_frame, text="Cancel", command=settings_window.destroy,
                 bg="#95a5a6", fg="white", width=15).pack(side=tk.LEFT, padx=5)
    
    def quick_image_to_pdf(self):
        folder = filedialog.askdirectory(title="Select folder with images")
        if folder:
            self.file_path.set(folder)
            self.target_format.set("PDF")
            self.convert_file()
    
    def quick_audio_batch(self):
        folder = filedialog.askdirectory(title="Select folder with audio files")
        if folder:
            self.file_path.set(folder)
            self.target_format.set("MP3")
            self.batch_convert()
    
    def quick_slideshow(self):
        folder = filedialog.askdirectory(title="Select folder with images for slideshow")
        if folder:
            self.file_path.set(folder)
            self.target_format.set("PPTX")
            self.ppt_layout.set("Multiple Images - One per Slide")
            self.convert_file()

    def quick_pdf_to_ppt(self):
        file = filedialog.askopenfilename(
            title="Select PDF file",
            filetypes=[("PDF files", "*.pdf"), ("All files", "*.*")]
        )
        if file:
            self.file_path.set(file)
            self.target_format.set("PPTX")
            self.convert_file()
    
    def browse_file(self):
        file_types = [
            ("All supported files", "*.jpg *.jpeg *.png *.bmp *.gif *.webp *.pdf *.docx *.doc *.txt *.pptx *.ppt *.mp3 *.wav *.ogg *.flac"),
            ("Image files", "*.jpg *.jpeg *.png *.bmp *.gif *.webp *.tiff *.ico"),
            ("PDF files", "*.pdf"),
            ("Document files", "*.docx *.doc *.txt *.rtf"),
            ("Presentation files", "*.pptx *.ppt"),
            ("Audio files", "*.mp3 *.wav *.ogg *.flac *.m4a *.aac"),
            ("All files", "*.*")
        ]
        filename = filedialog.askopenfilename(filetypes=file_types)
        if filename:
            self.file_path.set(filename)
            self.update_file_info(filename)
            self.update_conversion_options(filename)
    
    def browse_folder(self):
        folder = filedialog.askdirectory()
        if folder:
            self.file_path.set(folder)
            self.update_folder_info(folder)
    
    def update_conversion_options(self, filename):
        file_ext = Path(filename).suffix.lower()
        target_formats = []
        
        for category, info in self.supported_formats.items():
            if file_ext in info['extensions']:
                target_formats = info['convert_to']
                break
        
        if target_formats:
            self.format_combo['values'] = target_formats
            self.format_combo.set(target_formats[0])
            
            # Show/hide options based on target format
            current_format = self.target_format.get().upper()
            self.toggle_ppt_options(current_format == 'PPTX' and file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp'])
            self.toggle_pdf_ppt_options(file_ext == '.pdf' and current_format == 'PPTX')
            self.toggle_quality_option(current_format in ['JPEG', 'JPG', 'PNG', 'WEBP', 'TIFF'])
        else:
            self.format_combo.set('')
            self.format_combo['values'] = []
            self.log_message(f"Warning: No conversion options for {file_ext} files")
    
    def update_file_info(self, filename):
        try:
            file_stats = os.stat(filename)
            size_kb = file_stats.st_size / 1024
            file_ext = Path(filename).suffix.upper()
            modified_time = datetime.fromtimestamp(file_stats.st_mtime).strftime("%Y-%m-%d %H:%M:%S")
            
            info = f"📄 File: {os.path.basename(filename)}\n"
            info += f"📦 Size: {size_kb:.1f} KB | Type: {file_ext}\n"
            info += f"📅 Modified: {modified_time}\n"
            info += f"📁 Location: {os.path.dirname(filename)}\n"
            
            # Additional info based on file type
            if file_ext.lower() in ['.jpg', '.jpeg', '.png', '.bmp', '.gif', '.webp', '.tiff']:
                try:
                    with Image.open(filename) as img:
                        width, height = img.size
                        mode = img.mode
                        info += f"📐 Dimensions: {width}x{height} | Color Mode: {mode}\n"
                except:
                    info += "📐 Dimensions: Unable to read image properties\n"
            
            elif file_ext.lower() in ['.pptx', '.ppt']:
                info += "🎯 Type: PowerPoint Presentation\n"
                if PPTX_SUPPORT:
                    try:
                        from pptx import Presentation
                        prs = Presentation(filename)
                        info += f"📊 Slides: {len(prs.slides)} | "
                        info += "Format: PowerPoint\n"
                    except:
                        info += "📊 Slides: Unable to read presentation\n"
            
            elif file_ext.lower() == '.pdf' and FITZ_SUPPORT:
                try:
                    doc = fitz.open(filename)
                    info += f"📑 Pages: {len(doc)} | "
                    info += "Format: PDF Document\n"
                    doc.close()
                except:
                    info += "📑 Pages: Unable to read PDF\n"
            
            elif file_ext.lower() in ['.mp3', '.wav', '.ogg', '.flac'] and AUDIO_SUPPORT:
                try:
                    audio = AudioSegment.from_file(filename)
                    duration = len(audio) / 1000  # Convert to seconds
                    info += f"⏱️ Duration: {duration:.1f}s | "
                    info += f"Channels: {audio.channels} | Sample Rate: {audio.frame_rate}Hz\n"
                except:
                    info += "⏱️ Audio info: Unable to read audio properties\n"
            
            self.info_text.config(state=tk.NORMAL)
            self.info_text.delete(1.0, tk.END)
            self.info_text.insert(tk.END, info)
            self.info_text.config(state=tk.DISABLED)
            self.status.config(text=f"Selected: {os.path.basename(filename)}")
            
        except Exception as e:
            self.status.config(text=f"Error reading file info: {str(e)}")
    
    def update_folder_info(self, folder):
        try:
            files = []
            file_count_by_type = {}
            
            for category, info in self.supported_formats.items():
                for ext in info['extensions']:
                    found_files = list(Path(folder).glob(f"*{ext}")) + list(Path(folder).glob(f"*{ext.upper()}"))
                    files.extend(found_files)
                    if found_files:
                        file_count_by_type[category] = file_count_by_type.get(category, 0) + len(found_files)
            
            unique_files = list(set(files))
            
            info = f"📁 Folder: {os.path.basename(folder)}\n"
            info += f"📊 Total supported files: {len(unique_files)}\n"
            info += f"📍 Location: {folder}\n\n"
            info += "📋 File types found:\n"
            
            for file_type, count in file_count_by_type.items():
                icon = self.supported_formats[file_type].get('icon', '📄')
                info += f"  {icon} {file_type}: {count} files\n"
            
            info += f"\n🎯 First few files:\n"
            for i, file_path in enumerate(unique_files[:6]):
                info += f"  • {file_path.name}\n"
            
            if len(unique_files) > 6:
                info += f"  ... and {len(unique_files) - 6} more files"
            
            self.info_text.config(state=tk.NORMAL)
            self.info_text.delete(1.0, tk.END)
            self.info_text.insert(tk.END, info)
            self.info_text.config(state=tk.DISABLED)
            self.status.config(text=f"Selected folder: {os.path.basename(folder)}")
            
        except Exception as e:
            self.status.config(text=f"Error reading folder info: {str(e)}")
    
    def log_message(self, message):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_text.config(state=tk.NORMAL)
        self.log_text.insert(tk.END, f"[{timestamp}] {message}\n")
        self.log_text.see(tk.END)
        self.log_text.config(state=tk.DISABLED)
        self.root.update()
    
    def update_stats(self):
        self.stats_text.config(state=tk.NORMAL)
        self.stats_text.delete(1.0, tk.END)
        
        stats_info = f"📈 Conversion Statistics\n\n"
        stats_info += f"📊 Total files processed: {self.conversion_stats['total_files']}\n"
        stats_info += f"✅ Successful conversions: {self.conversion_stats['successful']}\n"
        stats_info += f"❌ Failed conversions: {self.conversion_stats['failed']}\n"
        
        if self.conversion_stats['total_files'] > 0:
            success_rate = (self.conversion_stats['successful'] / self.conversion_stats['total_files']) * 100
            stats_info += f"📈 Success rate: {success_rate:.1f}%\n"
        
        if self.conversion_stats['start_time']:
            stats_info += f"⏰ Tracking since: {self.conversion_stats['start_time']}\n"
        else:
            stats_info += f"⏰ No conversions yet\n"
        
        self.stats_text.insert(tk.END, stats_info)
        self.stats_text.config(state=tk.DISABLED)
    
    def get_output_path(self, input_path, target_format):
        """Determine output path based on user selection"""
        input_path = Path(input_path)
        base_name = f"{input_path.stem}_converted.{target_format.lower()}"
        
        if self.output_option.get() == "Same folder":
            return str(input_path.parent / base_name)
        elif self.output_option.get() == "Desktop":
            desktop = Path.home() / "Desktop"
            return str(desktop / base_name)
        else:  # Custom folder
            folder = filedialog.askdirectory(title="Select output folder")
            if folder:
                return str(Path(folder) / base_name)
            else:
                return str(input_path.parent / base_name)
    
    def convert_image(self, input_path, output_path, target_format):
        """Convert image files to various formats"""
        try:
            if target_format.upper() == 'PPTX':
                return self.image_to_pptx(input_path, output_path)
            elif target_format.upper() == 'PDF/A':
                # Convert to standard PDF first
                with Image.open(input_path) as img:
                    if img.mode in ('RGBA', 'P'):
                        img = img.convert('RGB')
                    img.save(output_path.replace('.pdfa', '.pdf'), 'PDF', resolution=100.0)
                return True, "Converted to PDF (PDF/A requires specialized library)"
            else:
                with Image.open(input_path) as img:
                    if target_format.lower() in ['jpeg', 'jpg']:
                        if img.mode in ('RGBA', 'P'):
                            img = img.convert('RGB')
                        img.save(output_path, quality=self.quality.get())
                    elif target_format.upper() == 'PDF':
                        if img.mode in ('RGBA', 'P'):
                            img = img.convert('RGB')
                        img.save(output_path, 'PDF', resolution=100.0)
                    else:
                        img.save(output_path)
                return True, "Success"
        except Exception as e:
            return False, str(e)
    
    def convert_pdf(self, input_path, output_path, target_format):
        """Convert PDF files"""
        try:
            if not FITZ_SUPPORT:
                return False, "PDF conversion requires PyMuPDF (fitz)"
            
            if target_format.upper() in ['PNG', 'JPG']:
                return self.convert_pdf_to_image(input_path, output_path, target_format)
            elif target_format.upper() == 'TXT':
                return self.convert_pdf_to_text(input_path, output_path)
            elif target_format.upper() == 'PPTX':
                return self.pdf_to_pptx_advanced(input_path, output_path)
            elif target_format.upper() == 'PDF/A':
                # Simple PDF to PDF/A conversion (basic implementation)
                doc = fitz.open(input_path)
                doc.save(output_path, garbage=4, deflate=True, clean=True)
                doc.close()
                return True, "Converted to optimized PDF"
            else:
                return False, f"Unsupported PDF conversion: {target_format}"
                
        except Exception as e:
            return False, str(e)
    
    def convert_pdf_to_text(self, input_path, output_path):
        """Convert PDF to text"""
        try:
            doc = fitz.open(input_path)
            text = ""
            for page in doc:
                text += page.get_text() + "\n"
            doc.close()
            
            with open(output_path, 'w', encoding='utf-8') as file:
                file.write(text)
            return True, "Success"
        except Exception as e:
            return False, str(e)
    
    def convert_pdf_to_image(self, input_path, output_path, target_format):
        """Convert PDF to images"""
        try:
            doc = fitz.open(input_path)
            output_files = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap()
                
                if len(doc) == 1:
                    # Single page - use original output path
                    page_output = output_path
                else:
                    # Multiple pages - add page number
                    name, ext = os.path.splitext(output_path)
                    page_output = f"{name}_page_{page_num+1}{ext}"
                
                if target_format.upper() == 'JPG':
                    pix.save(page_output)
                else:
                    pix.save(page_output)
                
                output_files.append(page_output)
            
            doc.close()
            return True, f"Created {len(output_files)} pages"
        except Exception as e:
            return False, str(e)
    
    def pdf_to_pptx_advanced(self, input_path, output_path):
        """Advanced PDF to PowerPoint conversion with multiple strategies"""
        try:
            if not PPTX_SUPPORT:
                return False, "PPTX conversion requires python-pptx"
            if not FITZ_SUPPORT:
                return False, "PDF to PPTX requires PyMuPDF (fitz)"
            
            doc = fitz.open(input_path)
            total_pages = len(doc)
            
            # Create a new presentation
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            
            conversion_mode = self.pdf_ppt_options.get()
            temp_files = []
            
            try:
                for page_num in range(total_pages):
                    page = doc.load_page(page_num)
                    
                    if conversion_mode == "Extract text content":
                        # Strategy 1: Extract text and create text-based slides
                        success = self._create_text_slide(prs, page, page_num, blank_slide_layout)
                        if not success:
                            # Fallback to image-based conversion
                            success = self._create_image_slide(prs, page, page_num, blank_slide_layout, temp_files)
                    
                    elif conversion_mode == "High quality images":
                        # Strategy 2: High-quality image conversion
                        success = self._create_high_quality_image_slide(prs, page, page_num, blank_slide_layout, temp_files)
                    
                    else:
                        # Strategy 3: Standard image conversion (default)
                        success = self._create_image_slide(prs, page, page_num, blank_slide_layout, temp_files)
                    
                    if not success:
                        self.log_message(f"Warning: Failed to convert page {page_num + 1}")
                
                # Save the presentation
                prs.save(output_path)
                result_message = f"Converted {total_pages} PDF pages to PowerPoint"
                
                if conversion_mode == "Extract text content":
                    result_message += " (text-based)"
                elif conversion_mode == "High quality images":
                    result_message += " (high-quality images)"
                else:
                    result_message += " (standard images)"
                
                return True, result_message
                
            finally:
                # Clean up temporary files
                for temp_file in temp_files:
                    try:
                        if os.path.exists(temp_file):
                            os.remove(temp_file)
                    except:
                        pass
                doc.close()
                
        except Exception as e:
            return False, f"PDF to PPTX conversion error: {str(e)}"
    
    def _create_text_slide(self, prs, page, page_num, blank_slide_layout):
        """Create a slide from extracted PDF text"""
        try:
            text = page.get_text()
            if not text.strip():
                return False  # No text found, fallback to image
            
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
            title_frame = title_shape.text_frame
            title_frame.text = f"Page {page_num + 1}"
            title_frame.paragraphs[0].font.size = Pt(24)
            title_frame.paragraphs[0].font.bold = True
            
            # Add content
            content_shape = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
            content_frame = content_shape.text_frame
            content_frame.text = text[:2000]  # Limit text length
            content_frame.paragraphs[0].font.size = Pt(12)
            
            return True
        except:
            return False
    
    def _create_image_slide(self, prs, page, page_num, blank_slide_layout, temp_files):
        """Create a slide from PDF page image (standard quality)"""
        try:
            # Render page as image
            mat = fitz.Matrix(1.5, 1.5)  # Standard resolution
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # Save temporary image
            temp_img_path = f"temp_pdf_page_{page_num}.png"
            with open(temp_img_path, "wb") as f:
                f.write(img_data)
            temp_files.append(temp_img_path)
            
            # Create slide and add image
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Calculate image dimensions to fit slide
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            img_ratio = pix.width / pix.height
            
            if img_ratio > (slide_width / slide_height):
                width = slide_width
                height = width / img_ratio
            else:
                height = slide_height
                width = height * img_ratio
            
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2
            
            slide.shapes.add_picture(temp_img_path, left, top, width, height)
            
            return True
        except:
            return False
    
    def _create_high_quality_image_slide(self, prs, page, page_num, blank_slide_layout, temp_files):
        """Create a slide from PDF page image (high quality)"""
        try:
            # Render page as high-quality image
            mat = fitz.Matrix(2.0, 2.0)  # Higher resolution
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # Save temporary image
            temp_img_path = f"temp_pdf_page_hq_{page_num}.png"
            with open(temp_img_path, "wb") as f:
                f.write(img_data)
            temp_files.append(temp_img_path)
            
            # Create slide and add image
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Calculate image dimensions to fit slide
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            img_ratio = pix.width / pix.height
            
            if img_ratio > (slide_width / slide_height):
                width = slide_width
                height = width / img_ratio
            else:
                height = slide_height
                width = height * img_ratio
            
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2
            
            slide.shapes.add_picture(temp_img_path, left, top, width, height)
            
            return True
        except:
            return False
    
    def image_to_pptx(self, input_path, output_path):
        """Convert image to PowerPoint presentation"""
        try:
            if not PPTX_SUPPORT:
                return False, "PPTX conversion requires python-pptx"
            
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            
            slide = prs.slides.add_slide(blank_slide_layout)
            
            with Image.open(input_path) as img:
                img_width, img_height = img.size
                img_ratio = img_width / img_height
            
            slide_width = Inches(10)
            slide_height = Inches(7.5)
            
            if img_ratio > (slide_width / slide_height):
                width = slide_width
                height = width / img_ratio
            else:
                height = slide_height
                width = height * img_ratio
            
            left = (slide_width - width) / 2
            top = (slide_height - height) / 2
            
            slide.shapes.add_picture(input_path, left, top, width, height)
            
            prs.save(output_path)
            return True, "Created PowerPoint with 1 slide"
            
        except Exception as e:
            return False, f"Image to PPTX error: {str(e)}"
    
    def images_to_pptx_batch(self, image_paths, output_path):
        """Convert multiple images to a single PowerPoint presentation"""
        try:
            if not PPTX_SUPPORT:
                return False, "PPTX conversion requires python-pptx"
            
            from pptx import Presentation
            from pptx.util import Inches
            
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]
            
            layout = self.ppt_layout.get()
            
            if layout == "Multiple Images - One per Slide":
                for img_path in image_paths:
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    with Image.open(img_path) as img:
                        img_width, img_height = img.size
                        img_ratio = img_width / img_height
                    
                    slide_width = Inches(10)
                    slide_height = Inches(7.5)
                    
                    if img_ratio > (slide_width / slide_height):
                        width = slide_width
                        height = width / img_ratio
                    else:
                        height = slide_height
                        width = height * img_ratio
                    
                    left = (slide_width - width) / 2
                    top = (slide_height - height) / 2
                    
                    slide.shapes.add_picture(str(img_path), left, top, width, height)
            
            prs.save(output_path)
            return True, f"Created PowerPoint with {len(image_paths)} slides"
            
        except Exception as e:
            return False, f"Batch images to PPTX error: {str(e)}"
    
    def images_to_pdf_batch(self, image_paths, output_path):
        """Convert multiple images to a single PDF"""
        try:
            images = []
            for p in image_paths:
                img = Image.open(p)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                images.append(img)
            
            if images:
                images[0].save(output_path, "PDF", resolution=100.0, save_all=True, append_images=images[1:])
                return True, f"Created PDF with {len(images)} pages"
            return False, "No images"
        except Exception as e:
            return False, str(e)
    
    def convert_presentation(self, input_path, output_path, target_format):
        """Convert PowerPoint presentations to various formats"""
        try:
            if not PPTX_SUPPORT:
                return False, "PowerPoint conversion requires python-pptx"
            
            file_ext = Path(input_path).suffix.lower()
            
            if target_format.upper() == 'PDF':
                return self.ppt_to_pdf(input_path, output_path)
            elif target_format.upper() in ['PNG', 'JPG']:
                return self.ppt_to_images(input_path, output_path, target_format)
            elif target_format.upper() == 'TXT':
                return self.ppt_to_text(input_path, output_path)
            elif target_format.upper() == 'PPTX':
                if file_ext == '.pptx':
                    prs = Presentation(input_path)
                    prs.save(output_path)
                    return True, "Converted presentation format"
                elif file_ext == '.ppt' and COMTYPES_SUPPORT:
                    import comtypes.client
                    powerpoint = comtypes.client.CreateObject('Powerpoint.Application')
                    powerpoint.Visible = False
                    abs_input = os.path.abspath(input_path)
                    abs_output = os.path.abspath(output_path)
                    pres = powerpoint.Presentations.Open(abs_input)
                    pres.SaveAs(abs_output, 11)  # 11 = ppSaveAsOpenXMLPresentation
                    pres.Close()
                    powerpoint.Quit()
                    return True, "Converted to PPTX"
                else:
                    return False, "Cannot convert to PPTX"
            else:
                return False, f"Unsupported presentation conversion: {target_format}"
                
        except Exception as e:
            return False, f"Presentation conversion error: {str(e)}"
    
    def ppt_to_pdf(self, input_path, output_path):
        """Convert PowerPoint to PDF"""
        try:
            if COMTYPES_SUPPORT:
                import comtypes.client
                powerpoint = comtypes.client.CreateObject('Powerpoint.Application')
                powerpoint.Visible = False
                abs_input = os.path.abspath(input_path)
                abs_output = os.path.abspath(output_path)
                pres = powerpoint.Presentations.Open(abs_input)
                pres.SaveAs(abs_output, 32)  # 32 = ppSaveAsPDF
                pres.Close()
                powerpoint.Quit()
                return True, "Success"
            else:
                # Fallback to text-based
                prs = Presentation(input_path)
                return self._create_pdf_from_ppt_text(prs, output_path)
                
        except Exception as e:
            return False, f"PPT to PDF error: {str(e)}"
    
    def _create_pdf_from_ppt_text(self, prs, output_path):
        """Create a simple PDF from PowerPoint text content"""
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            c = canvas.Canvas(output_path, pagesize=letter)
            width, height = letter
            
            y_position = height - 50
            c.setFont("Helvetica", 16)
            c.drawString(50, y_position, f"Presentation: {len(prs.slides)} Slides")
            y_position -= 30
            
            for i, slide in enumerate(prs.slides):
                if y_position < 100:
                    c.showPage()
                    y_position = height - 50
                    c.setFont("Helvetica", 12)
                
                c.setFont("Helvetica-Bold", 14)
                c.drawString(50, y_position, f"Slide {i+1}:")
                y_position -= 20
                c.setFont("Helvetica", 10)
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        lines = shape.text.split('\n')
                        for line in lines:
                            if line.strip():
                                if y_position < 50:
                                    c.showPage()
                                    y_position = height - 50
                                    c.setFont("Helvetica", 10)
                                c.drawString(70, y_position, line.strip())
                                y_position -= 15
                
                y_position -= 10
            
            c.save()
            return True, f"Created PDF with {len(prs.slides)} slides content"
        except Exception as e:
            return False, f"PDF creation error: {str(e)}"
    
    def ppt_to_images(self, input_path, output_path, target_format):
        """Convert PowerPoint to images"""
        try:
            if not COMTYPES_SUPPORT:
                return False, "COM not available for PPT rendering"
            import comtypes.client
            powerpoint = comtypes.client.CreateObject('Powerpoint.Application')
            powerpoint.Visible = True  # May need to be visible for export
            pres = powerpoint.Presentations.Open(os.path.abspath(input_path))
            base_name, _ = os.path.splitext(output_path)
            export_folder = base_name + "_slides"
            os.makedirs(export_folder, exist_ok=True)
            pres.Export(export_folder, target_format.upper())
            pres.Close()
            powerpoint.Quit()
            return True, f"Exported slides to {export_folder}"
        except Exception as e:
            return False, str(e)
    
    def ppt_to_text(self, input_path, output_path):
        """Extract text from PowerPoint presentation"""
        try:
            if not PPTX_SUPPORT:
                return False, "PPT to text requires python-pptx"
            
            prs = Presentation(input_path)
            text_content = []
            
            for i, slide in enumerate(prs.slides):
                text_content.append(f"--- Slide {i+1} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                text_content.append("")  # Empty line between slides
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(text_content))
            
            return True, f"Extracted text from {len(prs.slides)} slides"
            
        except Exception as e:
            return False, f"PPT to text error: {str(e)}"
    
    def convert_document(self, input_path, output_path, target_format):
        """Convert document files"""
        file_ext = Path(input_path).suffix.lower()
        if target_format.upper() == 'TXT':
            return self.document_to_text(input_path, output_path)
        elif target_format.upper() == 'PDF':
            return self.document_to_pdf(input_path, output_path)
        elif target_format.upper() == 'DOCX':
            return self.document_to_docx(input_path, output_path)
        else:
            return False, "Unsupported"
    
    def document_to_text(self, input_path, output_path):
        """Convert document to text"""
        try:
            file_ext = Path(input_path).suffix.lower()
            if file_ext == '.txt':
                with open(input_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                return True, "Success"
            elif file_ext == '.docx' and DOCX_SUPPORT:
                from docx import Document
                doc = Document(input_path)
                text = '\n'.join([para.text for para in doc.paragraphs])
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                return True, "Success"
            elif COMTYPES_SUPPORT and file_ext in ('.doc', '.rtf', '.docx'):
                import comtypes.client
                word = comtypes.client.CreateObject('Word.Application')
                word.Visible = False
                doc = word.Documents.Open(os.path.abspath(input_path))
                doc.SaveAs(os.path.abspath(output_path), FileFormat=0)  # 0 = txt
                doc.Close()
                word.Quit()
                return True, "Success"
            else:
                return False, "Cannot extract text"
        except Exception as e:
            return False, str(e)
    
    def document_to_pdf(self, input_path, output_path):
        """Convert document to PDF"""
        try:
            file_ext = Path(input_path).suffix.lower()
            if COMTYPES_SUPPORT:
                import comtypes.client
                word = comtypes.client.CreateObject('Word.Application')
                word.Visible = False
                abs_input = os.path.abspath(input_path)
                abs_output = os.path.abspath(output_path)
                doc = word.Documents.Open(abs_input)
                doc.SaveAs(abs_output, FileFormat=17)  # 17 = wdFormatPDF
                doc.Close()
                word.Quit()
                return True, "Success"
            elif file_ext == '.txt' and REPORTLAB_SUPPORT:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                y = height - 50
                c.setFont("Helvetica", 12)
                with open(input_path, 'r', encoding='utf-8') as f:
                    for line in f:
                        if y < 50:
                            c.showPage()
                            y = height - 50
                            c.setFont("Helvetica", 12)
                        c.drawString(50, y, line.strip())
                        y -= 15
                c.save()
                return True, "Success"
            elif file_ext == '.docx' and DOCX_SUPPORT and REPORTLAB_SUPPORT:
                from docx import Document
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                doc = Document(input_path)
                text = '\n'.join([para.text for para in doc.paragraphs])
                c = canvas.Canvas(output_path, pagesize=letter)
                width, height = letter
                y = height - 50
                c.setFont("Helvetica", 12)
                for line in text.split('\n'):
                    if y < 50:
                        c.showPage()
                        y = height - 50
                        c.setFont("Helvetica", 12)
                    c.drawString(50, y, line.strip())
                    y -= 15
                c.save()
                return True, "Success"
            else:
                return False, "Cannot convert to PDF"
        except Exception as e:
            return False, str(e)
    
    def document_to_docx(self, input_path, output_path):
        """Convert document to DOCX"""
        try:
            if not DOCX_SUPPORT:
                return False, "python-docx not available"
            file_ext = Path(input_path).suffix.lower()
            from docx import Document
            if file_ext == '.docx':
                doc = Document(input_path)
                doc.save(output_path)
                return True, "Success"
            elif file_ext == '.txt':
                doc = Document()
                with open(input_path, 'r', encoding='utf-8') as f:
                    for line in f:
                        doc.add_paragraph(line.strip())
                doc.save(output_path)
                return True, "Success"
            elif COMTYPES_SUPPORT and file_ext in ('.doc', '.rtf'):
                import comtypes.client
                word = comtypes.client.CreateObject('Word.Application')
                word.Visible = False
                abs_input = os.path.abspath(input_path)
                abs_output = os.path.abspath(output_path)
                doc = word.Documents.Open(abs_input)
                doc.SaveAs(abs_output, FileFormat=16)  # 16 = wdFormatXMLDocument
                doc.Close()
                word.Quit()
                return True, "Success"
            else:
                return False, "Cannot convert to DOCX"
        except Exception as e:
            return False, str(e)
    
    def convert_audio(self, input_path, output_path, target_format):
        """Convert audio files"""
        try:
            audio = AudioSegment.from_file(input_path)
            audio.export(output_path, format=target_format.lower())
            return True, "Success"
        except Exception as e:
            return False, str(e)
    
    def convert_file(self):
        # Initialize stats if first conversion
        if self.conversion_stats['start_time'] is None:
            self.conversion_stats['start_time'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        
        input_file = self.file_path.get()
        if not input_file:
            messagebox.showerror("Error", "Please select a file")
            return
        
        if not os.path.exists(input_file):
            messagebox.showerror("Error", "File does not exist")
            return
        
        # Run conversion in thread to prevent UI freezing
        thread = threading.Thread(target=self._convert_file_thread, daemon=True)
        thread.start()
    
    def _convert_file_thread(self):
        """Threaded file conversion to prevent UI freezing"""
        try:
            input_file = self.file_path.get()
            
            if os.path.isdir(input_file):
                # Collect files
                all_files = []
                for category, info in self.supported_formats.items():
                    for ext in info['extensions']:
                        all_files.extend(list(Path(input_file).glob(f"*{ext}")))
                        all_files.extend(list(Path(input_file).glob(f"*{ext.upper()}")))
                unique_files = list(set(all_files))
                
                # Handle folder conversion for images to PPTX or PDF
                if self.target_format.get().upper() == 'PPTX':
                    image_files = [f for f in unique_files if f.suffix.lower() in self.supported_formats.get('Images', {}).get('extensions', [])]
                    if image_files:
                        output_file = self.get_output_path(input_file, "pptx")
                        success, message = self.images_to_pptx_batch(image_files, output_file)
                        
                        self.conversion_stats['total_files'] += len(image_files)
                        if success:
                            self.conversion_stats['successful'] += len(image_files)
                            self.log_message(f"✓ Success: Created PPTX from {len(image_files)} images")
                            self.root.after(0, lambda: messagebox.showinfo("Success", 
                                          f"Presentation created successfully!\n\n"
                                          f"Images processed: {len(image_files)}\n"
                                          f"Output: {os.path.basename(output_file)}\n"
                                          f"Layout: {self.ppt_layout.get()}"))
                            self.status.config(text="Presentation created successfully!")
                        else:
                            self.conversion_stats['failed'] += len(image_files)
                            self.log_message(f"✗ Failed: {message}")
                            self.root.after(0, lambda: messagebox.showerror("Error", f"Conversion failed:\n{message}"))
                            self.status.config(text="Conversion failed")
                        self.update_stats()
                        return
                elif self.target_format.get().upper() == 'PDF':
                    image_files = [f for f in unique_files if f.suffix.lower() in self.supported_formats.get('Images', {}).get('extensions', [])]
                    if image_files:
                        output_file = self.get_output_path(input_file, "pdf")
                        success, message = self.images_to_pdf_batch(image_files, output_file)
                        
                        self.conversion_stats['total_files'] += len(image_files)
                        if success:
                            self.conversion_stats['successful'] += len(image_files)
                            self.log_message(f"✓ Success: Created PDF from {len(image_files)} images")
                            self.root.after(0, lambda: messagebox.showinfo("Success", 
                                          f"PDF created successfully!\n\n"
                                          f"Images processed: {len(image_files)}\n"
                                          f"Output: {os.path.basename(output_file)}"))
                            self.status.config(text="PDF created successfully!")
                        else:
                            self.conversion_stats['failed'] += len(image_files)
                            self.log_message(f"✗ Failed: {message}")
                            self.root.after(0, lambda: messagebox.showerror("Error", f"Conversion failed:\n{message}"))
                            self.status.config(text="Conversion failed")
                        self.update_stats()
                        return
                else:
                    self.batch_convert()
                    return
            
            self.progress.start()
            self.status.config(text="Converting...")
            self.log_message(f"Starting conversion: {os.path.basename(input_file)}")
            
            # Create output filename
            target_format = self.target_format.get().lower()
            output_file = self.get_output_path(input_file, target_format)
            
            file_ext = Path(input_file).suffix.lower()
            success = False
            message = ""
            
            # Route to appropriate converter
            if file_ext in self.supported_formats.get('Images', {}).get('extensions', []):
                success, message = self.convert_image(input_file, output_file, target_format)
            elif file_ext in self.supported_formats.get('PDF', {}).get('extensions', []):
                success, message = self.convert_pdf(input_file, output_file, target_format)
            elif file_ext in self.supported_formats.get('Documents', {}).get('extensions', []):
                success, message = self.convert_document(input_file, output_file, target_format)
            elif file_ext in self.supported_formats.get('Presentations', {}).get('extensions', []):
                success, message = self.convert_presentation(input_file, output_file, target_format)
            elif file_ext in self.supported_formats.get('Audio', {}).get('extensions', []):
                success, message = self.convert_audio(input_file, output_file, target_format)
            else:
                message = f"Unsupported file type: {file_ext}"
            
            self.conversion_stats['total_files'] += 1
            if success:
                self.conversion_stats['successful'] += 1
                self.log_message(f"✓ Success: {os.path.basename(output_file)}")
                self.root.after(0, lambda: messagebox.showinfo("Success", 
                              f"File converted successfully!\n\n"
                              f"Original: {os.path.basename(input_file)}\n"
                              f"Converted: {os.path.basename(output_file)}\n"
                              f"Format: {self.target_format.get()}"))
                self.status.config(text="Conversion completed successfully!")
            else:
                self.conversion_stats['failed'] += 1
                self.log_message(f"✗ Failed: {message}")
                self.root.after(0, lambda: messagebox.showerror("Error", f"Conversion failed:\n{message}"))
                self.status.config(text="Conversion failed")
            
            self.update_stats()
            
        except Exception as e:
            error_msg = f"Unexpected error: {str(e)}"
            self.log_message(f"✗ Error: {error_msg}")
            self.root.after(0, lambda: messagebox.showerror("Error", error_msg))
            self.status.config(text="Conversion failed")
        finally:
            self.progress.stop()
            self.progress_label.config(text="0%")
    
    def batch_convert(self):
        folder = self.file_path.get()
        if not folder or not os.path.isdir(folder):
            folder = filedialog.askdirectory(title="Select folder with files")
            if not folder:
                return
        
        # Run batch conversion in thread
        thread = threading.Thread(target=self._batch_convert_thread, args=(folder,), daemon=True)
        thread.start()
    
    def _batch_convert_thread(self, folder):
        """Threaded batch conversion"""
        try:
            # Collect all supported files
            all_files = []
            for category, info in self.supported_formats.items():
                for ext in info['extensions']:
                    all_files.extend(list(Path(folder).glob(f"*{ext}")))
                    all_files.extend(list(Path(folder).glob(f"*{ext.upper()}")))
            
            unique_files = list(set(all_files))
            
            if not unique_files:
                self.root.after(0, lambda: messagebox.showinfo("No Files", "No supported files found in the selected folder."))
                return
            
            target_format = self.target_format.get().lower()
            converted_count = 0
            total_files = len(unique_files)
            
            self.progress.config(maximum=total_files, value=0)
            self.status.config(text=f"Converting {total_files} files...")
            self.log_message(f"Starting batch conversion of {total_files} files to {target_format}")
            
            for i, file_path in enumerate(unique_files):
                try:
                    input_path = str(file_path)
                    output_path = self.get_output_path(input_path, target_format)
                    
                    file_ext = file_path.suffix.lower()
                    success = False
                    message = ""
                    
                    if file_ext in self.supported_formats.get('Images', {}).get('extensions', []):
                        success, message = self.convert_image(input_path, output_path, target_format)
                    elif file_ext in self.supported_formats.get('PDF', {}).get('extensions', []):
                        success, message = self.convert_pdf(input_path, output_path, target_format)
                    elif file_ext in self.supported_formats.get('Documents', {}).get('extensions', []):
                        success, message = self.convert_document(input_path, output_path, target_format)
                    elif file_ext in self.supported_formats.get('Presentations', {}).get('extensions', []):
                        success, message = self.convert_presentation(input_path, output_path, target_format)
                    elif file_ext in self.supported_formats.get('Audio', {}).get('extensions', []):
                        success, message = self.convert_audio(input_path, output_path, target_format)
                    
                    if success:
                        converted_count += 1
                        self.log_message(f"✓ {file_path.name}")
                    else:
                        self.log_message(f"✗ {file_path.name}: {message}")
                    
                    progress_value = i + 1
                    progress_percent = int((progress_value / total_files) * 100)
                    self.progress['value'] = progress_value
                    self.progress_label.config(text=f"{progress_percent}%")
                    self.status.config(text=f"Progress: {i+1}/{total_files}")
                    
                except Exception as e:
                    self.log_message(f"✗ Error with {file_path.name}: {str(e)}")
            
            self.conversion_stats['total_files'] += total_files
            self.conversion_stats['successful'] += converted_count
            self.conversion_stats['failed'] += (total_files - converted_count)
            
            self.root.after(0, lambda: messagebox.showinfo("Batch Complete", 
                              f"Batch conversion completed!\n\n"
                              f"Total files processed: {total_files}\n"
                              f"Successfully converted: {converted_count}\n"
                              f"Failed: {total_files - converted_count}"))
            self.status.config(text=f"Batch completed: {converted_count}/{total_files} files")
            self.update_stats()
            
        except Exception as e:
            self.root.after(0, lambda: messagebox.showerror("Error", f"Batch conversion failed:\n{str(e)}"))
            self.status.config(text="Batch conversion failed")
        finally:
            self.progress['value'] = 0
            self.progress_label.config(text="0%")
    
    def run(self):
        self.root.mainloop()

if __name__ == "__main__":
    app = UniversalConverter()
    app.run()
